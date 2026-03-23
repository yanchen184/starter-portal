"""
產生提案 PPT：團隊協作與知識管理平台
完全對齊 提案內容.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── 色彩定義 ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
ACCENT_ORANGE = RGBColor(0xE8, 0x8D, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
SUBTITLE_GRAY = RGBColor(0x99, 0x99, 0x99)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
CARD_BG = RGBColor(0x2A, 0x2A, 0x45)
CARD_BG2 = RGBColor(0x23, 0x23, 0x3D)
FONT = 'Microsoft JhengHei'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return box


def add_multiline(slide, left, top, width, height, lines, size=18,
                  color=LIGHT_GRAY, line_spacing=1.3):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(size * 0.4)
    return box


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    return slide


def add_title(slide, text, accent_width=Inches(2.0)):
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             text, size=36, color=WHITE, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.15), accent_width, Pt(4), ACCENT_BLUE)


def add_table(slide, headers, rows, top, col_widths, left=Inches(0.8)):
    """用 shape 模擬表格，更好看"""
    col_x = [left]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    # Header
    for j, header in enumerate(headers):
        add_rect(slide, col_x[j], top, col_widths[j], Inches(0.5), ACCENT_BLUE)
        add_text(slide, col_x[j] + Inches(0.08), top + Inches(0.02),
                 col_widths[j] - Inches(0.16), Inches(0.45),
                 header, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Rows
    row_h = Inches(0.55)
    for i, row in enumerate(rows):
        y = top + Inches(0.5) + i * row_h
        bg = CARD_BG if i % 2 == 0 else CARD_BG2
        for j, cell in enumerate(row):
            add_rect(slide, col_x[j], y, col_widths[j], row_h, bg)
            add_text(slide, col_x[j] + Inches(0.08), y + Inches(0.02),
                     col_widths[j] - Inches(0.16), row_h - Inches(0.04),
                     cell, size=12, color=LIGHT_GRAY)


# ════════════════════════════════════════
# Slide 1: 封面
# ════════════════════════════════════════
slide = new_slide()
add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
         '團隊協作與知識管理 提案',
         size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
         '打造高效團隊的兩大基礎建設',
         size=24, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(4.2), Inches(2.3), Pt(3), ACCENT_BLUE)
add_text(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
         '2026 Q2  |  內部提案',
         size=18, color=SUBTITLE_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════
# Slide 2: 核心理念
# ════════════════════════════════════════
slide = new_slide()
add_title(slide, '核心理念', Inches(1.5))

add_text(slide, Inches(2), Inches(2.2), Inches(9), Inches(1.0),
         '「好的團隊，不應該因為少了誰就無法運作。」',
         size=36, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_multiline(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1.2), [
    '知識如果只存在某個人的腦袋裡，',
    '那個人請假、換組、離職，工作就會停擺。',
    '這不是個人能力的問題，而是組織缺乏知識管理的系統性風險。',
], size=20, color=LIGHT_GRAY)

goals = [
    '需求規格、技術設計、決策脈絡，全部有文字紀錄',
    '新人看文件就能上手，不用靠老人口耳相傳',
    '不再有「這件事只有某某知道」的狀況',
    '人員異動不影響團隊的知識傳承與運作效率',
]
for i, g in enumerate(goals):
    add_text(slide, Inches(2.5), Inches(5.0) + i * Inches(0.55), Inches(8), Inches(0.5),
             f'✓  {g}', size=18, color=GREEN)

# ════════════════════════════════════════
# Slide 3: 現況痛點
# ════════════════════════════════════════
slide = new_slide()
add_title(slide, '一、現況痛點', Inches(2.0))

# 左欄：文檔問題
add_rect(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.0), CARD_BG)
add_text(slide, Inches(0.9), Inches(1.95), Inches(5.2), Inches(0.5),
         '1.1  文檔與規格書', size=24, color=ACCENT_ORANGE, bold=True)
doc_pain = [
    '幾乎沒有 Spec 文件，需求靠口頭傳遞',
    '技術文件散落各 Git Repo，難以搜尋',
    '新人 onboarding 沒有參考資料，養成時間長',
    '同樣的問題重複被問、重複踩坑',
    'PM 與工程師之間資訊落差大',
]
add_multiline(slide, Inches(0.9), Inches(2.7), Inches(5.2), Inches(3.5),
              [f'•  {p}' for p in doc_pain], size=17)

# 右欄：溝通問題
add_rect(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.0), CARD_BG)
add_text(slide, Inches(7.2), Inches(1.95), Inches(5.2), Inches(0.5),
         '1.2  溝通工具', size=24, color=ACCENT_ORANGE, bold=True)
comm_pain = [
    'LINE 訊息容易被洗掉、無法追溯',
    '無法針對單一議題進行討論（Thread）',
    '新成員加入看不到歷史訊息脈絡',
    '無法整合開發工具通知（Git、CI/CD）',
    '重要決策散落在聊天記錄中',
]
add_multiline(slide, Inches(7.2), Inches(2.7), Inches(5.2), Inches(3.5),
              [f'•  {p}' for p in comm_pain], size=17)

# ════════════════════════════════════════
# Slide 4: 隱性成本
# ════════════════════════════════════════
slide = new_slide()
add_title(slide, '1.3  隱性成本', Inches(2.0))
add_rect(slide, Inches(0.8), Inches(1.15), Inches(2.0), Pt(4), RED)

costs = [
    ('溝通成本', '每人每天花 30~60 分鐘在「找資料」和「問人」\n80 人團隊 = 每月損失 200+ 人時', '200+\n人時/月'),
    ('重複犯錯', '沒有知識沉澱，同樣的坑不同人重複踩\n修一次 bug 可能要 2~4 小時', '2~4 hr\n/ 次'),
    ('新人養成慢', '沒有文檔支撐，onboarding 靠師徒制\n新人從入職到產出的時間拉長 2~3 倍', '2~3 倍\n更久'),
    ('決策不透明', '重要脈絡隨人員異動而消失\n沒人知道「當初為什麼這樣設計」', '無法\n追溯'),
]

for i, (title, desc, metric) in enumerate(costs):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + col * Inches(6.3)
    y = Inches(1.8) + row * Inches(2.6)
    add_rect(slide, x, y, Inches(5.8), Inches(2.2), CARD_BG)
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.8), Inches(0.5),
             title, size=22, color=RED, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.75), Inches(3.8), Inches(1.3),
             desc, size=16, color=LIGHT_GRAY)
    add_text(slide, x + Inches(4.0), y + Inches(0.3), Inches(1.5), Inches(1.5),
             metric, size=28, color=RED, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════
# Slide 5: 知識管理平台需求
# ════════════════════════════════════════
slide = new_slide()
add_title(slide, '二、知識管理平台需求', Inches(3.0))

# P0
add_rect(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.2), CARD_BG)
add_text(slide, Inches(0.9), Inches(1.85), Inches(5.2), Inches(0.5),
         'P0 核心需求', size=22, color=RED, bold=True)

add_text(slide, Inches(0.9), Inches(2.5), Inches(5.2), Inches(0.4),
         '權限控管', size=18, color=ACCENT_BLUE, bold=True)
add_multiline(slide, Inches(0.9), Inches(2.9), Inches(5.2), Inches(1.8), [
    '•  依專案/部門分權',
    '•  依角色分權（PM / RD / QA / 主管）',
    '•  「全公司可見」與「專案內可見」兩種層級',
    '•  由管理者指定存取權限',
], size=15)

add_text(slide, Inches(0.9), Inches(4.5), Inches(5.2), Inches(0.4),
         '內容模板', size=18, color=ACCENT_BLUE, bold=True)
add_multiline(slide, Inches(0.9), Inches(4.9), Inches(5.2), Inches(2.0), [
    '•  PRD（需求規格書）— PM 撰寫',
    '•  TDD（技術設計文件）— 工程師撰寫',
    '•  API Spec — 工程師撰寫',
    '•  Onboarding 指南 — 每個專案一份',
    '•  Meeting Notes — 會議決策紀錄',
], size=15)

# P1 + P2
add_rect(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(2.4), CARD_BG)
add_text(slide, Inches(7.2), Inches(1.85), Inches(5.2), Inches(0.5),
         'P1 重要需求', size=22, color=ACCENT_ORANGE, bold=True)
add_multiline(slide, Inches(7.2), Inches(2.5), Inches(5.2), Inches(1.5), [
    '•  全文搜尋 — 文件多了之後找得到才有用',
    '•  版本歷史 — 誰改了什麼、能回溯',
    '•  協作編輯 — 多人同時編輯',
], size=15)

add_rect(slide, Inches(6.9), Inches(4.3), Inches(5.8), Inches(2.6), CARD_BG)
add_text(slide, Inches(7.2), Inches(4.45), Inches(5.2), Inches(0.5),
         'P2 加分需求', size=22, color=GREEN, bold=True)
add_multiline(slide, Inches(7.2), Inches(5.1), Inches(5.2), Inches(1.5), [
    '•  結構化組織 — 樹狀目錄/空間/標籤',
    '•  整合能力 — Git、Slack、CI/CD 串接',
    '•  匯出功能 — PDF/Word 匯出',
], size=15)

# ════════════════════════════════════════
# Slide 6: 知識管理平台比較
# ════════════════════════════════════════
slide = new_slide()
add_title(slide, '三、知識管理平台選項比較', Inches(3.5))

headers = ['平台', '類型', '權限控管', '內容模板', '搜尋', '版本歷史', '協作編輯', 'PM 友善', '費用/月']
rows_data = [
    ['Confluence', 'SaaS', '強', '強（必填欄位）', '強', '完整比對', '有', '中', '~$580'],
    ['Notion', 'SaaS', '中', '強（Database）', '強', '有', '有', '高', '~$800'],
    ['Outline', '自架/SaaS', '有', '有', '有', '有', '有', '高', '免費(自架)'],
    ['BookStack', '自架', '有', '有', '有', '有', '有限', '高', '免費(自架)'],
    ['GitBook', 'SaaS', '有', '有', '有', '有', '有', '中', '~$500'],
    ['VitePress', '靜態站', '無', '弱', '有', '靠Git', '無', '低', '免費'],
    ['Slite', 'SaaS', '有', '有', '有', '有', '有', '高', '~$640'],
    ['Coda', 'SaaS', '有', '強', '有', '有', '有', '高', '~$720'],
    ['MediaWiki', '自架', '有', '有', '有', '有', '有限', '低', '免費(自架)'],
]

cw = [Inches(1.3), Inches(1.1), Inches(1.1), Inches(1.6), Inches(0.8), Inches(1.2), Inches(1.1), Inches(1.0), Inches(1.3)]
add_table(slide, headers, rows_data, Inches(1.6), cw, left=Inches(0.5))

add_text(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
         'VitePress 適合作為公開技術文件站，但缺乏權限控管、協作編輯，不適合作為全團隊知識管理平台',
         size=14, color=SUBTITLE_GRAY)

# ════════════════════════════════════════
# Slide 7: 溝通平台需求 + 比較
# ════════════════════════════════════════
slide = new_slide()
add_title(slide, '四 & 五、溝通平台需求與比較', Inches(4.0))

# 需求
add_rect(slide, Inches(0.6), Inches(1.7), Inches(5.0), Inches(5.2), CARD_BG)
add_text(slide, Inches(0.9), Inches(1.85), Inches(4.5), Inches(0.5),
         '核心需求', size=22, color=ACCENT_BLUE, bold=True)
add_multiline(slide, Inches(0.9), Inches(2.5), Inches(4.5), Inches(2.0), [
    '•  訊息永久保留（不可有時間限制）',
    '•  頻道/群組管理（依專案/團隊分頻道）',
    '•  Thread 討論（針對單一議題回覆）',
    '•  全文搜尋歷史訊息',
], size=16)

add_text(slide, Inches(0.9), Inches(4.5), Inches(4.5), Inches(0.5),
         '重要需求', size=22, color=ACCENT_ORANGE, bold=True)
add_multiline(slide, Inches(0.9), Inches(5.1), Inches(4.5), Inches(1.5), [
    '•  整合開發工具（GitHub、CI/CD）',
    '•  檔案分享與保留',
    '•  新成員加入可看到歷史訊息',
], size=16)

# 比較表
headers2 = ['平台', '訊息保留', 'Thread', '整合', '費用/月', '備註']
rows2 = [
    ['Slack Pro', '無限', '有', '最豐富', '~$590', '業界標配'],
    ['Slack Free', '90 天', '有', '豐富', '免費', '不適合長期'],
    ['Mattermost', '無限(自架)', '有', '不錯', '伺服器成本', '需維運'],
    ['Discord', '無限', '有', '有限', '免費', '不夠正式'],
    ['Google Chat', '需GWS', '有', 'Google系', '~$960', '需GWS'],
    ['MS Teams', '無限', '有', 'Office系', '~$480', '需M365'],
]

cw2 = [Inches(1.4), Inches(1.3), Inches(0.8), Inches(1.1), Inches(1.3), Inches(1.5)]
add_table(slide, headers2, rows2, Inches(1.7), cw2, left=Inches(6.0))

# ════════════════════════════════════════
# Slide 8: 方案 A
# ════════════════════════════════════════
slide = new_slide()
add_title(slide, '六、組合方案', Inches(2.0))

# 三個方案卡片
plans = [
    {
        'title': '方案 A — 零成本自架',
        'color': GREEN,
        'items': [
            ('知識管理', 'Outline（自架）'),
            ('溝通平台', 'Mattermost（自架）'),
            ('技術文件', 'starter-portal (VitePress)'),
            ('月費', '僅伺服器成本'),
        ],
        'pros': '零授權費\n資料完全自控',
        'cons': '需要維運人力\n生態系較小',
    },
    {
        'title': '方案 B — 快速導入',
        'color': ACCENT_BLUE,
        'items': [
            ('知識管理', 'Notion'),
            ('溝通平台', 'Slack Pro'),
            ('技術文件', 'starter-portal + Notion 同步'),
            ('月費', '~$1,400 USD'),
        ],
        'pros': '上手最快\nPM 友善度最高\n零維運',
        'cons': '費用較高\n權限控管不如\nConfluence',
    },
    {
        'title': '方案 C — 企業級',
        'color': ACCENT_ORANGE,
        'items': [
            ('知識管理', 'Confluence'),
            ('溝通平台', 'Slack Pro'),
            ('技術文件', 'starter-portal + Confluence 同步'),
            ('月費', '~$1,170 USD'),
        ],
        'pros': '權限最完整\n企業功能最強\n可整合 Jira',
        'cons': '學習曲線稍高\n介面較重',
    },
]

for i, plan in enumerate(plans):
    x = Inches(0.4) + i * Inches(4.3)
    add_rect(slide, x, Inches(1.7), Inches(4.0), Inches(5.3), CARD_BG)
    add_rect(slide, x, Inches(1.7), Inches(4.0), Pt(5), plan['color'])

    add_text(slide, x + Inches(0.2), Inches(1.9), Inches(3.6), Inches(0.5),
             plan['title'], size=20, color=plan['color'], bold=True)

    for j, (key, val) in enumerate(plan['items']):
        y = Inches(2.6) + j * Inches(0.45)
        add_text(slide, x + Inches(0.2), y, Inches(1.2), Inches(0.4),
                 key, size=14, color=SUBTITLE_GRAY)
        add_text(slide, x + Inches(1.4), y, Inches(2.4), Inches(0.4),
                 val, size=14, color=WHITE)

    add_text(slide, x + Inches(0.2), Inches(4.6), Inches(1.6), Inches(0.3),
             '優點', size=14, color=GREEN, bold=True)
    add_text(slide, x + Inches(0.2), Inches(4.9), Inches(3.6), Inches(1.0),
             plan['pros'], size=14, color=LIGHT_GRAY)

    add_text(slide, x + Inches(2.0), Inches(4.6), Inches(1.6), Inches(0.3),
             '缺點', size=14, color=RED, bold=True)
    add_text(slide, x + Inches(2.0), Inches(4.9), Inches(1.8), Inches(1.0),
             plan['cons'], size=14, color=LIGHT_GRAY)

# ════════════════════════════════════════
# Slide 9: 推動計畫
# ════════════════════════════════════════
slide = new_slide()
add_title(slide, '七、分階段推動計畫', Inches(3.0))
add_rect(slide, Inches(0.8), Inches(1.15), Inches(3.0), Pt(4), GREEN)

phases = [
    ('Phase 1 — 試行期', '第 1~2 個月', ACCENT_BLUE, [
        '選定 1 個專案團隊 (5~10 人) 先行試用',
        '建立溝通平台 Workspace，設定基本頻道',
        '在知識平台上建立文件模板',
        'starter-portal 技術文件 POC 持續運作',
    ]),
    ('Phase 2 — 擴展期', '第 3~4 個月', ACCENT_ORANGE, [
        '根據試行回饋調整流程與模板',
        '擴展到 2~3 個團隊',
        '建立文件撰寫規範 & 最佳實踐',
        '整合 Git / CI 通知到溝通平台',
    ]),
    ('Phase 3 — 全面導入', '第 5~6 個月', GREEN, [
        '全團隊導入',
        '建立 onboarding 文件流程',
        '定期 review 文件品質',
        '評估是否需要調整或升級方案',
    ]),
]

for i, (title, period, color, items) in enumerate(phases):
    x = Inches(0.4) + i * Inches(4.3)
    add_rect(slide, x, Inches(1.7), Inches(4.0), Inches(5.2), CARD_BG)
    add_rect(slide, x, Inches(1.7), Inches(4.0), Pt(5), color)

    add_text(slide, x + Inches(0.2), Inches(1.9), Inches(3.6), Inches(0.5),
             title, size=20, color=color, bold=True)
    add_text(slide, x + Inches(0.2), Inches(2.4), Inches(3.6), Inches(0.4),
             period, size=16, color=SUBTITLE_GRAY)

    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.2), Inches(3.1) + j * Inches(0.6),
                 Inches(3.6), Inches(0.55),
                 f'•  {item}', size=15, color=LIGHT_GRAY)

# ════════════════════════════════════════
# Slide 10: 預期效益
# ════════════════════════════════════════
slide = new_slide()
add_title(slide, '八、預期效益', Inches(2.0))

benefits = [
    ('減少溝通成本', '每人每天節省 30 分鐘找資料時間\n80 人 × 0.5hr × 22天', '880\n人時/月'),
    ('加速新人上手', '有文檔支撐的 onboarding\n大幅縮短養成時間', '50%\n更快'),
    ('減少重複錯誤', '知識沉澱避免重複踩坑\n減少重複問題發生', '60%\n減少'),
    ('決策可追溯', '所有討論和決策都有紀錄\n人員異動不影響知識傳承', '100%\n可追溯'),
]

for i, (title, desc, metric) in enumerate(benefits):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + col * Inches(6.3)
    y = Inches(1.8) + row * Inches(2.6)
    add_rect(slide, x, y, Inches(5.8), Inches(2.2), CARD_BG)
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.8), Inches(0.5),
             title, size=22, color=GREEN, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.75), Inches(3.8), Inches(1.3),
             desc, size=16, color=LIGHT_GRAY)
    add_text(slide, x + Inches(4.0), y + Inches(0.3), Inches(1.5), Inches(1.5),
             metric, size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════
# Slide 11: 下一步行動
# ════════════════════════════════════════
slide = new_slide()
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.0),
         '九、下一步行動', size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(2.5), Inches(2.3), Pt(3), ACCENT_BLUE)

next_steps = [
    '1.  決定知識管理平台選型',
    '2.  決定溝通平台選型',
    '3.  選定試行團隊，開始 Phase 1',
    '4.  指定文件負責人，建立初始模板',
]

for i, step in enumerate(next_steps):
    add_text(slide, Inches(3.0), Inches(3.2) + i * Inches(0.7), Inches(7.5), Inches(0.6),
             step, size=24, color=LIGHT_GRAY)

add_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.6),
         '小步快跑，先求有、再求好',
         size=28, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════
# 儲存
# ════════════════════════════════════════
output_path = '/Users/yanchen/workspace/my-projects/starter-portal/團隊協作與知識管理提案.pptx'
prs.save(output_path)
print(f'PPT 已產生：{output_path}')
